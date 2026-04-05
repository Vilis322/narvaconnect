"""
Parse all Moodle files (PDF, DOCX, XLSX, PPTX, HTML, IPYNB, CSV, TXT)
into structured text, then generate training Q&A pairs.

Runs locally — no API tokens.

Usage:
    python data/scripts/parse_moodle.py
"""

import json
import re
from collections import defaultdict
from datetime import datetime
from pathlib import Path

# ============================================================
# File parsers — each returns extracted text
# ============================================================

def parse_pdf(path: str) -> str:
    """Extract text from PDF using pymupdf."""
    try:
        import fitz
        doc = fitz.open(path)
        text = ""
        for page in doc:
            text += page.get_text() + "\n"
        doc.close()
        return text.strip()
    except Exception as e:
        return f"[PDF parse error: {e}]"


def parse_docx(path: str) -> str:
    """Extract text from DOCX."""
    try:
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        return f"[DOCX parse error: {e}]"


def parse_pptx(path: str) -> str:
    """Extract text from PPTX slides."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append(f"[Slide {i}]\n" + "\n".join(slide_text))
        return "\n\n".join(text_parts)
    except Exception as e:
        return f"[PPTX parse error: {e}]"


def parse_xlsx(path: str) -> str:
    """Extract text from XLSX using pandas."""
    try:
        import pandas as pd
        dfs = pd.read_excel(path, sheet_name=None)
        text_parts = []
        for sheet_name, df in dfs.items():
            text_parts.append(f"[Sheet: {sheet_name}]")
            text_parts.append(df.to_string(index=False, max_rows=50))
        return "\n\n".join(text_parts)
    except Exception as e:
        return f"[XLSX parse error: {e}]"


def parse_csv(path: str) -> str:
    """Extract text from CSV."""
    try:
        import pandas as pd
        df = pd.read_csv(path, nrows=50)
        return f"Columns: {', '.join(df.columns)}\n\n{df.to_string(index=False, max_rows=20)}"
    except Exception as e:
        return f"[CSV parse error: {e}]"


def parse_html(path: str) -> str:
    """Extract text from HTML (Moodle pages)."""
    try:
        with open(path, encoding="utf-8", errors="ignore") as f:
            html = f.read()
        # Simple tag stripping
        text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
        text = re.sub(r'<[^>]+>', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()
        # Decode HTML entities
        import html as html_lib
        text = html_lib.unescape(text)
        return text[:5000]  # Limit HTML text
    except Exception as e:
        return f"[HTML parse error: {e}]"


def parse_ipynb(path: str) -> str:
    """Extract code and markdown cells from Jupyter notebook."""
    try:
        with open(path, encoding="utf-8") as f:
            nb = json.load(f)
        text_parts = []
        for cell in nb.get("cells", []):
            cell_type = cell.get("cell_type", "")
            source = "".join(cell.get("source", []))
            if source.strip():
                prefix = "[Code]" if cell_type == "code" else "[Text]"
                text_parts.append(f"{prefix}\n{source.strip()}")
        return "\n\n".join(text_parts[:30])  # Limit cells
    except Exception as e:
        return f"[IPYNB parse error: {e}]"


def parse_txt(path: str) -> str:
    """Read plain text file."""
    try:
        with open(path, encoding="utf-8", errors="ignore") as f:
            return f.read()[:5000]
    except Exception as e:
        return f"[TXT parse error: {e}]"


# Parser registry
PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".doc": parse_docx,  # try same parser
    ".pptx": parse_pptx,
    ".xlsx": parse_xlsx,
    ".csv": parse_csv,
    ".html": parse_html,
    ".ipynb": parse_ipynb,
    ".txt": parse_txt,
    ".py": parse_txt,
    ".json": parse_txt,
}

# Skip these
SKIP_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".mp4", ".pkl", ".data", ""}

# Subject folder → readable name
SUBJECT_NAMES = {
    "analysis_and_design_of_is": "Analysis and Design of Information Systems",
    "intro_to_data_science": "Introduction to Data Science",
    "software_engineering": "Software Engineering",
    "software_testing": "Software Testing",
    "starting_a_business": "Starting a Business",
    "web_application_development": "Web Application Development",
    "practical_project_training": "Practical Project Training",
    "introductory_practice": "Introductory Practice",
}

UNIVERSITY_DIR = Path.home() / "Projects/study/university/year3/semester2"
OUTPUT_DIR = Path(__file__).resolve().parent.parent / "processed"


def detect_subject(file_path: Path) -> str:
    """Detect subject from file path."""
    parts = file_path.parts
    for i, part in enumerate(parts):
        if part == "moodle" and i > 0:
            folder = parts[i - 1]
            return SUBJECT_NAMES.get(folder, folder.replace("_", " ").title())
    return "Unknown"


def generate_qa_from_document(subject: str, filename: str, text: str) -> list[dict]:
    """Generate Q&A training pairs from a parsed document."""
    pairs = []

    if not text or len(text) < 20:
        return pairs
    # Skip only actual error messages, not structured content
    if text.startswith("[") and "parse error" in text.lower():
        return pairs

    # Truncate very long texts
    text_truncated = text[:3000]

    # General "what is this file about" pair
    pairs.append({
        "prompt": f"What is the file '{filename}' about in {subject}?",
        "completion": f"The file '{filename}' from {subject} contains: {text_truncated[:500]}...",
    })

    # If it looks like an assignment
    assignment_keywords = ["assignment", "task", "submit", "deadline", "ülesanne",
                          "homework", "kodutöö", "praktikum", "exercise"]
    if any(kw in text.lower() for kw in assignment_keywords):
        pairs.append({
            "prompt": f"What assignments are there in {subject}?",
            "completion": f"Based on '{filename}': {text_truncated[:800]}",
        })

    # If it's a lecture/presentation
    lecture_keywords = ["lecture", "loeng", "slide", "chapter", "topic"]
    if any(kw in text.lower() or kw in filename.lower() for kw in lecture_keywords):
        pairs.append({
            "prompt": f"What topics are covered in {subject} lectures?",
            "completion": f"From '{filename}': {text_truncated[:800]}",
        })

    # If it's a notebook (practical work)
    if filename.endswith(".ipynb"):
        pairs.append({
            "prompt": f"What practical exercises are in {subject}?",
            "completion": f"Notebook '{filename}' contains: {text_truncated[:800]}",
        })

    return pairs


def main():
    print("=" * 60)
    print("  NarvaConnect — Moodle File Parser")
    print("  PDF, DOCX, XLSX, PPTX, HTML, IPYNB, CSV, TXT")
    print("  Local parsing, no API tokens")
    print("=" * 60)

    # Find all moodle files
    all_files = []
    for moodle_dir in sorted(UNIVERSITY_DIR.glob("*/moodle")):
        for f in sorted(moodle_dir.rglob("*")):
            if f.is_file() and f.name != ".DS_Store":
                all_files.append(f)

    # Also check nested moodle dirs (like intro_to_data_science has sub-dirs)
    for moodle_dir in sorted(UNIVERSITY_DIR.glob("*/moodle/**/moodle")):
        for f in sorted(moodle_dir.rglob("*")):
            if f.is_file() and f.name != ".DS_Store":
                all_files.append(f)

    print(f"\nFound {len(all_files)} files to parse\n")

    parsed_docs = []
    training_pairs = []
    stats = defaultdict(int)

    for i, file_path in enumerate(all_files):
        ext = file_path.suffix.lower()
        subject = detect_subject(file_path)
        filename = file_path.name

        if ext in SKIP_EXTENSIONS:
            stats["skipped"] += 1
            continue

        parser = PARSERS.get(ext)
        if not parser:
            print(f"[{i+1}/{len(all_files)}] {filename} — no parser for {ext}")
            stats["no_parser"] += 1
            continue

        print(f"[{i+1}/{len(all_files)}] {subject} / {filename}", end=" ... ")
        text = parser(str(file_path))

        is_error = not text or (text.startswith("[") and "parse error" in text.lower())
        is_empty = text and len(text.strip()) < 20

        if is_error or is_empty:
            print(f"FAILED ({text[:50] if text else 'empty'})")
            stats["failed"] += 1
            continue

        char_count = len(text)
        print(f"{char_count} chars")
        stats["parsed"] += 1

        parsed_docs.append({
            "subject": subject,
            "filename": filename,
            "extension": ext,
            "text_length": char_count,
            "text_preview": text[:200],
        })

        # Generate Q&A pairs
        qa_pairs = generate_qa_from_document(subject, filename, text)
        training_pairs.extend(qa_pairs)
        stats["qa_pairs"] += len(qa_pairs)

    # Write parsed docs index
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    docs_path = OUTPUT_DIR / "moodle_parsed.json"
    with open(docs_path, "w") as f:
        json.dump({
            "documents": parsed_docs,
            "parsed_at": datetime.now().isoformat(),
            "total_files": len(all_files),
            "stats": dict(stats),
        }, f, indent=2, ensure_ascii=False)

    # Append to existing training.jsonl (or create new)
    training_path = OUTPUT_DIR / "moodle_training.jsonl"
    with open(training_path, "w") as f:
        for pair in training_pairs:
            f.write(json.dumps(pair, ensure_ascii=False) + "\n")

    print(f"\n{'=' * 60}")
    print("  Done!")
    print(f"  Files found: {len(all_files)}")
    print(f"  Parsed: {stats['parsed']}")
    print(f"  Skipped (binary): {stats['skipped']}")
    print(f"  Failed: {stats['failed']}")
    print(f"  No parser: {stats['no_parser']}")
    print(f"  Training pairs: {stats['qa_pairs']}")
    print(f"  Output: {docs_path}")
    print(f"  Training: {training_path}")
    print(f"{'=' * 60}")
    print("\nTo merge with ois2 training data:")
    print("  cat data/processed/training.jsonl data/processed/moodle_training.jsonl > data/processed/combined_training.jsonl")


if __name__ == "__main__":
    main()
